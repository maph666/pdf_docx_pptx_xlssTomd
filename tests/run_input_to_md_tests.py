import sys
from pathlib import Path

# ensure project root is importable
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from converter import convert_file

OUT = Path(__file__).parent / 'out'
OUT.mkdir(exist_ok=True)

# create sample docx
try:
    import docx
    d = docx.Document()
    d.add_heading('Documento de Prueba', level=1)
    d.add_paragraph('Este es un párrafo de ejemplo.')
    p = d.add_paragraph()
    p.add_run('• Punto 1')
    p = d.add_paragraph()
    p.add_run('• Punto 2')
    table = d.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = 'Col1'
    table.rows[0].cells[1].text = 'Col2'
    table.rows[1].cells[0].text = 'a'
    table.rows[1].cells[1].text = '1'
    docx_path = OUT / 'sample.docx'
    d.save(str(docx_path))
except Exception as e:
    print('SKIP docx sample creation:', e)
    docx_path = None

# create sample pptx
try:
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Diapositiva 1'
    body = slide.shapes.placeholders[1].text_frame
    body.text = 'Punto 1\n- Subpunto'
    pptx_path = OUT / 'sample.pptx'
    prs.save(str(pptx_path))
except Exception as e:
    print('SKIP pptx sample creation:', e)
    pptx_path = None

# create sample xlsx
try:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = 'Hoja1'
    ws.append(['Col1', 'Col2'])
    ws.append(['a', 1])
    xlsx_path = OUT / 'sample.xlsx'
    wb.save(str(xlsx_path))
except Exception as e:
    print('SKIP xlsx sample creation:', e)
    xlsx_path = None


def run_conversion(path, suffix='md'):
    out = OUT / (path.stem + '.' + suffix)
    try:
        convert_file(str(path), str(out))
        print(f'OK: converted {path} -> {out}')
        print(out.read_text(encoding='utf-8')[:400])
    except Exception as e:
        print(f'ERROR converting {path}:', e)


if docx_path:
    run_conversion(docx_path)
if pptx_path:
    run_conversion(pptx_path)
if xlsx_path:
    run_conversion(xlsx_path)
