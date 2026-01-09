from markitdown import MarkItDown
from pathlib import Path
import sys


def convert_file(input_path: str, output_path: str):
    """Convert input file (PDF/DOCX/PPTX/XLSX) to Markdown and write to output_path.

    Raises a RuntimeError if required dependencies are missing or the input type is unsupported.
    """
    input_path = str(input_path)
    output_path = str(output_path)

    suffix = Path(input_path).suffix.lower()

    if suffix == '.pdf':
        md = MarkItDown()
        result = md.convert(input_path)
        content = None
        for attr in ("markdown", "md", "text_content", "text"):
            content = getattr(result, attr, None)
            if content:
                break
        if content is None:
            content = str(result)
    elif suffix == '.docx':
        content = docx_to_md(input_path)
    elif suffix == '.pptx':
        content = pptx_to_md(input_path)
    elif suffix in ('.xlsx', '.xls'):
        content = xlsx_to_md(input_path)
    else:
        raise RuntimeError(f"Unsupported input file type: {suffix}")

    Path(output_path).write_text(content, encoding='utf-8')
    return output_path


# --- extractors: file -> markdown text ---

def docx_to_md(docx_path: str) -> str:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError("Word conversion requires 'python-docx'. Install it and retry.") from exc

    doc = docx.Document(docx_path)
    out_lines = []

    # paragraphs
    for p in doc.paragraphs:
        text = p.text or ''
        text = text.strip()
        if not text:
            out_lines.append('')
            continue
        style = ''
        try:
            style = (p.style.name or '').lower()
        except Exception:
            style = ''
        # headings
        if style.startswith('heading'):
            import re
            m = re.search(r"\d+", style)
            level = int(m.group()) if m else 2
            out_lines.append('#' * level + ' ' + text)
        # simple list detection
        elif style.startswith('list') or text.startswith(('•', '·', '-', '*')):
            # normalize bullets
            if text.startswith(('- ', '* ')):
                out_lines.append('- ' + text[2:].strip())
            else:
                out_lines.append('- ' + text.lstrip('•·').strip())
        else:
            out_lines.append(text)

    # tables (append after paragraphs)
    for table in doc.tables:
        rows = []
        for r in table.rows:
            rows.append([c.text.strip() for c in r.cells])
        if not rows:
            continue
        # header
        header = rows[0]
        out_lines.append('')
        out_lines.append('| ' + ' | '.join(header) + ' |')
        out_lines.append('| ' + ' | '.join(['---'] * len(header)) + ' |')
        for r in rows[1:]:
            out_lines.append('| ' + ' | '.join(r) + ' |')
        out_lines.append('')

    return '\n'.join(out_lines).strip() + '\n'


def pptx_to_md(pptx_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PowerPoint conversion requires 'python-pptx'. Install it and retry.") from exc

    prs = Presentation(pptx_path)
    out_lines = []

    for slide in prs.slides:
        title = None
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = None
        if title:
            out_lines.append('# ' + title)
        # extract text from shapes
        for shape in slide.shapes:
            if not hasattr(shape, 'text'):
                continue
            text = shape.text.strip()
            if not text:
                continue
            # split into paragraphs and preserve simple bullets/levels
            for para in text.splitlines():
                para = para.strip()
                if not para:
                    continue
                if para.startswith(('-', '*', '•')):
                    out_lines.append('- ' + para.lstrip('-*• ').strip())
                else:
                    out_lines.append(para)
        out_lines.append('')

    return '\n'.join(out_lines).strip() + '\n'


def xlsx_to_md(xlsx_path: str) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("Excel conversion requires 'openpyxl'. Install it and retry.") from exc

    wb = load_workbook(xlsx_path, data_only=True)
    out_lines = []

    for sheet in wb.worksheets:
        out_lines.append('# ' + sheet.title)
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            out_lines.append('')
            continue
        # treat first row as header
        header = [str(c) if c is not None else '' for c in rows[0]]
        out_lines.append('| ' + ' | '.join(header) + ' |')
        out_lines.append('| ' + ' | '.join(['---'] * len(header)) + ' |')
        for r in rows[1:]:
            out_lines.append('| ' + ' | '.join([str(c) if c is not None else '' for c in r]) + ' |')
        out_lines.append('')

    return '\n'.join(out_lines).strip() + '\n'
